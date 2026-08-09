"""Text extraction for PDF papers and PPTX presentations."""
from __future__ import annotations

import os
import re

from pypdf import PdfReader
from pptx import Presentation


class ExtractionError(Exception):
    pass


def _clean(text: str) -> str:
    text = text.replace("\x00", " ").replace("\u200b", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_pdf(path: str) -> dict:
    """Return per-page text plus full text for a PDF file."""
    pages = {}
    full_parts = []
    try:
        reader = PdfReader(path)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"Could not read PDF: {exc}") from exc

    for index, page in enumerate(reader.pages, start=1):
        try:
            raw = page.extract_text() or ""
        except Exception:  # some malformed pages raise
            raw = ""
        text = _clean(raw)
        if text:
            pages[index] = text
            full_parts.append(text)

    if not pages:
        raise ExtractionError("No extractable text found in the PDF (scanned pages are not supported).")

    return {
        "kind": "pdf",
        "num_pages": len(pages),
        "pages": pages,
        "full_text": _clean("\n\n".join(full_parts)),
    }


def _shape_text(shape) -> str:
    parts = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            run_text = "".join(run.text for run in para.runs) or para.text
            if run_text.strip():
                parts.append(run_text)
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        except Exception:  # noqa: BLE001
            pass
    if getattr(shape, "shape_type", None) is not None and str(shape.shape_type) == "PICTURE (13)":
        parts.append("[IMAGE]")
    return "\n".join(parts)


def extract_pptx(path: str) -> dict:
    """Return per-slide text plus full text for a PPTX file."""
    slides = {}
    slide_chars = []
    try:
        prs = Presentation(path)
    except Exception as exc:  # noqa: BLE001
        raise ExtractionError(f"Could not read PPTX: {exc}") from exc

    for index, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text.strip():
                parts.append(text)
        text = _clean("\n".join(parts))
        if not text:
            text = "[NO TEXT ON SLIDE]"
        slides[index] = text
        slide_chars.append(text)

    if not slides:
        raise ExtractionError("The PPTX does not contain any slides.")

    return {
        "kind": "pptx",
        "num_slides": len(slides),
        "slides": slides,
        "full_text": _clean("\n\n".join(slide_chars)),
    }


def extract_from_path(path: str) -> dict:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    raise ExtractionError(f"Unsupported file type: {ext} (expected .pdf or .pptx)")